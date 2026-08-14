import os
import io
import subprocess
import tempfile

import fitz
from docx import Document
from pptx import Presentation

import config


def _ocr_paddle(page):
    from paddleocr import PaddleOCR

    ocr = PaddleOCR(lang=config.OCR_LANGUAGE)
    pix = page.get_pixmap(dpi=300)

    tmp_path = os.path.join(tempfile.gettempdir(), "_ocr_page.png")
    with open(tmp_path, "wb") as f:
        f.write(pix.tobytes("png"))

    result = ocr.ocr(tmp_path)
    lines = []
    for block in result:
        if isinstance(block, dict) and "rec_texts" in block:
            lines.extend(block["rec_texts"])
            continue
        for line in block:
            lines.append(line[1][0])
    return "\n".join(lines)


def _ocr_tesseract(page):
    import os
    import pytesseract
    from PIL import Image

    if config.TESSERACT_CMD_PATH:
        pytesseract.pytesseract.tesseract_cmd = config.TESSERACT_CMD_PATH
    if config.TESSDATA_PREFIX:
        os.environ["TESSDATA_PREFIX"] = config.TESSDATA_PREFIX

    pix = page.get_pixmap(dpi=300)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    return pytesseract.image_to_string(img, lang=config.TESSERACT_LANGUAGE)


def ocr_page(page):
    try:
        return _ocr_paddle(page)
    except Exception:
        return _ocr_tesseract(page)


def load_pdf(file_path):
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if not text:
            text = ocr_page(page)
        pages.append({
            "page_number": i,
            "text": text,
            "source": os.path.basename(file_path),
        })
    doc.close()
    return pages


def load_docx(file_path):
    document = Document(file_path)
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [{"page_number": 1, "text": text, "source": os.path.basename(file_path)}]


def load_doc(file_path):
    with tempfile.TemporaryDirectory() as tmp_dir:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "txt", "--outdir", tmp_dir, file_path],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        txt_path = os.path.join(tmp_dir, base_name + ".txt")
        with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

    return [{"page_number": 1, "text": text, "source": os.path.basename(file_path)}]


def load_pptx(file_path):
    presentation = Presentation(file_path)
    pages = []
    for i, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        pages.append({
            "page_number": i,
            "text": "\n".join(texts),
            "source": os.path.basename(file_path),
        })
    return pages


def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [{"page_number": 1, "text": text, "source": os.path.basename(file_path)}]


LOADERS = {
    "pdf": load_pdf,
    "docx": load_docx,
    "doc": load_doc,
    "pptx": load_pptx,
    "ppt": load_pptx,
    "txt": load_txt,
}


def load_document(file_path):
    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    if ext not in config.SUPPORTED_FILE_TYPES:
        raise ValueError(f"Unsupported file type: {ext}")
    return LOADERS[ext](file_path)
